import os
import zipfile
import shutil
import tempfile
import xml.etree.ElementTree as ET
from pathlib import Path
from datetime import datetime
import json
import random
import re
from io import BytesIO
import math
from concurrent.futures import ThreadPoolExecutor, as_completed
import multiprocessing as mp
import logging

logging.basicConfig(level=logging.INFO, format='%(asctime)s | %(levelname)s | %(message)s', datefmt='%Y-%m-%d %H:%M:%S')
logger = logging.getLogger(__name__)

from pydub import AudioSegment
from pydub.effects import normalize
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


class ModernPresentationGenerator:
    def __init__(self, base_path: str):
        self.base_path = Path(base_path)
        
        if not self.base_path.exists():
            raise FileNotFoundError(f"Шаблон не найден: {self.base_path}")

        self.round_slide_ranges = {
            1: (5, 44),
            2: (48, 87),
            3: (91, 130)
        }
        
        self.service_slides = {1, 2, 3, 4, 45, 46, 47, 88, 89, 90, 131}
        
        self.buffer_ms = 5000
        self.default_ms = 35_000
        self.min_duration = 5_000
        self.max_duration = 120_000
        self.fade_duration = 3000
        self.fade_in_duration = 1000

        self.min_image_width = 750
        self.min_image_height = 750
        self.max_image_width = 1600
        self.max_image_height = 1200
        self.target_size_for_presentation = 800
        self.small_image_threshold = 800

        self._red_words_per_slide = {}
        self._artist_red_words_cache = {}
        self._round_track_counters = {1: 0, 2: 0, 3: 0}

    def _delete_slides_directly(self, pptx_path: Path, slides_to_delete: set) -> Path:
        try:
            logger.info(f"Удаляем слайды: {sorted(slides_to_delete)}")
            
            tmp_dir = Path(tempfile.mkdtemp(prefix="delete_slides_"))
            extract_dir = tmp_dir / "extracted"
            with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
                zip_ref.extractall(extract_dir)
            
            deleted_count = 0
            slides_dir = extract_dir / "ppt" / "slides"
            rels_dir = slides_dir / "_rels"
            
            for slide_num in slides_to_delete:
                slide_file = slides_dir / f"slide{slide_num}.xml"
                if slide_file.exists():
                    slide_file.unlink()
                    deleted_count += 1
                rel_file = rels_dir / f"slide{slide_num}.xml.rels"
                if rel_file.exists():
                    rel_file.unlink()
            
            logger.info(f"Удалено файлов слайдов: {deleted_count}")
            
            pres_rels_path = extract_dir / "ppt" / "_rels" / "presentation.xml.rels"
            if pres_rels_path.exists():
                tree = ET.parse(pres_rels_path)
                root = tree.getroot()
                ET.register_namespace('', 'http://schemas.openxmlformats.org/package/2006/relationships')
                
                relationships = root.findall('{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')
                rels_deleted = 0
                for rel in relationships:
                    target = rel.get('Target', '')
                    for slide_num in slides_to_delete:
                        if f"slides/slide{slide_num}.xml" in target:
                            root.remove(rel)
                            rels_deleted += 1
                            break
                
                tree.write(pres_rels_path, encoding='UTF-8', xml_declaration=True)
                logger.info(f"Удалено ссылок в presentation.xml.rels: {rels_deleted}")
            
            pres_path = extract_dir / "ppt" / "presentation.xml"
            if pres_path.exists():
                rels_tree = ET.parse(pres_rels_path)
                rels_root = rels_tree.getroot()
                
                remaining_rids = {}
                for rel in rels_root.findall('{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                    target = rel.get('Target', '')
                    rid = rel.get('Id', '')
                    match = re.search(r'slides/slide(\d+)\.xml', target)
                    if match and rid:
                        slide_num = int(match.group(1))
                        remaining_rids[rid] = slide_num
                
                tree = ET.parse(pres_path)
                root = tree.getroot()
                ET.register_namespace('p', 'http://schemas.openxmlformats.org/presentationml/2006/main')
                ET.register_namespace('r', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships')
                ET.register_namespace('a', 'http://schemas.openxmlformats.org/drawingml/2006/main')
                
                slide_id_list = root.find('{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst')
                if slide_id_list:
                    slide_ids = slide_id_list.findall('{http://schemas.openxmlformats.org/presentationml/2006/main}sldId')
                    slides_removed = 0
                    for slide_id in list(slide_ids):
                        rid = slide_id.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                        if rid and rid not in remaining_rids:
                            slide_id_list.remove(slide_id)
                            slides_removed += 1
                    
                    remaining_slides = slide_id_list.findall('{http://schemas.openxmlformats.org/presentationml/2006/main}sldId')
                    for idx, slide_id in enumerate(remaining_slides, 1):
                        slide_id.set('id', str(256 + idx))
                    
                    logger.info(f"Удалено slideId из presentation.xml: {slides_removed}")
                
                tree.write(pres_path, encoding='UTF-8', xml_declaration=True)
            
            new_pptx = tmp_dir / "presentation_without_slides.pptx"
            
            with zipfile.ZipFile(new_pptx, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for root, dirs, files in os.walk(extract_dir):
                    for file in files:
                        file_path = Path(root) / file
                        skip_file = False
                        for slide_num in slides_to_delete:
                            if f"slide{slide_num}.xml" in str(file_path) or f"slide{slide_num}.xml.rels" in str(file_path):
                                skip_file = True
                                break
                        
                        if not skip_file:
                            arcname = file_path.relative_to(extract_dir)
                            zipf.write(file_path, arcname)
            
            logger.info(f"Создан новый PPTX без удаленных слайдов: {new_pptx}")
            shutil.copy2(new_pptx, pptx_path)
            shutil.rmtree(tmp_dir, ignore_errors=True)
            
            logger.info(f"Слайды удалены из: {pptx_path}")
            return pptx_path
            
        except Exception as e:
            logger.error(f"Ошибка удаления слайдов: {e}")
            import traceback
            logger.error(traceback.format_exc())
            raise

    def _calculate_slides_to_delete(self, rounds_config: list, rounds_count: int):
        slides_to_delete = set()
        logger.info(f"Рассчет удаления слайдов для {rounds_count} раундов: {rounds_config}")
        
        if rounds_count == 1:
            slides_to_delete.update(range(45, 48))
            slides_to_delete.update(range(88, 91))
            logger.info("Режим 1 раунда: удаляем слайды 45-47 и 88-90")
        
        elif rounds_count == 2:
            slides_to_delete.update(range(88, 91))
            logger.info("Режим 2 раундов: удаляем слайды 88-90")
        
        elif rounds_count == 3:
            logger.info("Режим 3 раундов: все служебные слайды остаются")
        
        for round_idx in range(rounds_count):
            round_num = round_idx + 1
            tracks_in_round = rounds_config[round_idx] if round_idx < len(rounds_config) else 40
            start_slide, end_slide = self.round_slide_ranges[round_num]
            last_track_slide = start_slide + tracks_in_round - 1
            
            if last_track_slide < end_slide:
                for slide in range(last_track_slide + 1, end_slide + 1):
                    slides_to_delete.add(slide)
            
            logger.info(f"Раунд {round_num}: {tracks_in_round} треков, слайды {start_slide}-{last_track_slide}")
        
        for round_num in range(rounds_count + 1, 4):
            if round_num in self.round_slide_ranges:
                start_slide, end_slide = self.round_slide_ranges[round_num]
                for slide in range(start_slide, end_slide + 1):
                    slides_to_delete.add(slide)
                logger.info(f"Удаляем все слайды неактивного раунда {round_num}: {start_slide}-{end_slide}")
        
        slides_to_delete.discard(1)
        slides_to_delete.discard(2)
        slides_to_delete.discard(3)
        slides_to_delete.discard(4)
        slides_to_delete.discard(131)
        
        sorted_deletes = sorted(slides_to_delete)
        logger.info(f"Всего слайдов для удаления ({len(sorted_deletes)}): {sorted_deletes}")
        
        return slides_to_delete

    def _crop_transparent_borders(self, img: Image.Image) -> Image.Image:
        try:
            if img.mode != 'RGBA':
                img = img.convert('RGBA')
            
            alpha = img.getchannel('A')
            bbox = alpha.getbbox()
            
            if bbox:
                left, upper, right, lower = bbox
                padding_w = int((right - left) * 0.02)
                padding_h = int((lower - upper) * 0.02)
                left = max(0, left - padding_w)
                upper = max(0, upper - padding_h)
                right = min(img.width, right + padding_w)
                lower = min(img.height, lower + padding_h)
                img = img.crop((left, upper, right, lower))
            
            return img
            
        except Exception as e:
            logger.warning(f"Ошибка обрезки границ: {e}")
            return img

    def _normalize_image_size(self, img: Image.Image) -> Image.Image:
        try:
            original_width, original_height = img.size
            is_horizontal = original_width > original_height
            is_small = original_width < self.small_image_threshold and original_height < self.small_image_threshold
            
            if is_small:
                logger.info("Маленькое изображение - увеличиваем в 1.5 раза")
                
                if is_horizontal:
                    target_width = self.min_image_width
                    scale = target_width / original_width
                    target_height = int(original_height * scale)
                    
                    if target_height < self.min_image_height:
                        scale = self.min_image_height / original_height
                        target_width = int(original_width * scale)
                        target_height = self.min_image_height
                        
                        if target_width > self.max_image_width:
                            scale = self.max_image_width / target_width
                            target_width = self.max_image_width
                            target_height = int(target_height * scale)
                else:
                    target_height = self.min_image_height
                    scale = target_height / original_height
                    target_width = int(original_width * scale)
                    
                    if target_width < self.min_image_width:
                        scale = self.min_image_width / original_width
                        target_width = self.min_image_width
                        target_height = int(original_height * scale)
                        
                        if target_height > self.max_image_height:
                            scale = self.max_image_height / target_height
                            target_height = self.max_image_height
                            target_width = int(target_width * scale)
                
                img = img.resize((target_width, target_height), Image.LANCZOS)
                logger.info(f"Маленькое фото увеличено: {original_width}x{original_height} -> {target_width}x{target_height}")
                return img
            
            aspect_ratio = original_width / original_height
            
            if is_horizontal:
                scale_factor = 0.75
                target_width = min(self.max_image_width, int(original_width * scale_factor))
                target_height = int(target_width / aspect_ratio)
                
                if target_width < self.min_image_width:
                    target_width = self.min_image_width
                    target_height = int(target_width / aspect_ratio)
                if target_height < self.min_image_height:
                    target_height = self.min_image_height
                    target_width = int(target_height * aspect_ratio)
                    
                if target_width > self.max_image_width:
                    target_width = self.max_image_width
                    target_height = int(target_width / aspect_ratio)
                if target_height > self.max_image_height:
                    target_height = self.max_image_height
                    target_width = int(target_height * aspect_ratio)
            else:
                scale_factor = 0.85
                target_height = min(self.max_image_height, int(original_height * scale_factor))
                target_width = int(target_height * aspect_ratio)
                
                if target_height < self.min_image_height:
                    target_height = self.min_image_height
                    target_width = int(target_height * aspect_ratio)
                if target_width < self.min_image_width:
                    target_width = self.min_image_width
                    target_height = int(target_width / aspect_ratio)
                    
                if target_height > self.max_image_height:
                    target_height = self.max_image_height
                    target_width = int(target_height * aspect_ratio)
                if target_width > self.max_image_width:
                    target_width = self.max_image_width
                    target_height = int(target_width / aspect_ratio)
            
            if (target_width, target_height) != (original_width, original_height):
                img = img.resize((target_width, target_height), Image.LANCZOS)
            
            return img
            
        except Exception as e:
            logger.warning(f"Ошибка нормализации размера: {e}")
            return img

    def _normalize_audio_loudness(self, audio_segment):
        try:
            logger.info("Нормализация громкости аудио")
            return normalize(audio_segment, headroom=1.0)
        except Exception as e:
            logger.warning(f"Ошибка нормализации: {e}")
            return audio_segment

    def _validate_segment_duration(self, duration_seconds: float) -> int:
        min_seconds = self.min_duration / 1000
        max_seconds = self.max_duration / 1000
        validated_seconds = max(min_seconds, min(duration_seconds, max_seconds))
        return int(validated_seconds * 1000)

    def _get_image_orientation(self, image_path: str):
        try:
            with Image.open(image_path) as img:
                width, height = img.size
                orientation = "horizontal" if width > height else "vertical"
                return orientation
        except Exception as e:
            logger.warning(f"Ошибка определения ориентации: {e}")
            return "vertical"

    def _resize_image_to_max_height(self, img: Image.Image, max_height_px: int = 800, scale_factor: float = 1.34) -> Image.Image:
        w, h = img.size
        if h <= max_height_px:
            base_img = img
        else:
            ratio = max_height_px / h
            new_w = int(w * ratio)
            new_h = max_height_px
            base_img = img.resize((new_w, new_h), Image.LANCZOS)
        
        final_w = int(base_img.width * scale_factor)
        final_h = int(base_img.height * scale_factor)
        return base_img.resize((final_w, final_h), Image.LANCZOS)

    def _load_and_process_image_fast(self, image_path: str, make_bw: bool):
        path = Path(image_path)
        if not path.exists():
            logger.warning(f"Изображение не найдено: {image_path}")
            return None
        
        try:
            logger.info(f"Загрузка и обработка изображения: {image_path}")
            
            with Image.open(path) as img:
                if img.mode == "P":
                    img = img.convert("RGBA")
                elif img.mode == "L":
                    img = img.convert("RGBA")
                elif img.mode == "RGB":
                    img = img.convert("RGBA")
                
                img = self._crop_transparent_borders(img)
                img = self._normalize_image_size(img)
                
                if make_bw:
                    r, g, b, a = img.split()
                    grayscale = ImageOps.grayscale(img.convert("RGB"))
                    img = Image.merge("RGBA", (grayscale, grayscale, grayscale, a))
                
                img = self._resize_image_to_max_height(img, max_height_px=self.target_size_for_presentation, scale_factor=1.34)
                
                logger.info(f"Изображение обработано: {img.mode}, размер: {img.size}")
                return img.copy()

        except Exception as e:
            logger.error(f"Ошибка обработки изображения {image_path}: {e}")
            return None

    def _apply_fade_effects(self, audio_segment):
        try:
            fade_in = min(self.fade_in_duration, len(audio_segment) // 3)
            fade_out = min(self.fade_duration, len(audio_segment) - 1000)
            if fade_in > 0:
                audio_segment = audio_segment.fade_in(fade_in)
            if fade_out > 0:
                audio_segment = audio_segment.fade_out(fade_out)
            return audio_segment
        except Exception as e:
            logger.warning(f"Ошибка применения эффектов: {e}")
            return audio_segment

    def _find_track_file_fast(self, track_path: str):
        possible_paths = [
            Path(track_path),
            Path.cwd() / "downloads" / Path(track_path).name,
            Path.cwd() / "uploads" / Path(track_path).name
        ]
        for path in possible_paths:
            if path.exists():
                return path
        logger.warning(f"Файл трека не найден: {track_path}")
        return None

    def _load_tracks_from_json_fast(self):
        possible_paths = [
            Path.cwd() / "tracks.json",
            Path.cwd() / "Track_data.json", 
            Path.cwd() / "track_data.json",
        ]
        for path in possible_paths:
            if path.exists():
                try:
                    logger.info(f"Загрузка треков из: {path}")
                    with open(path, "r", encoding="utf-8") as f:
                        data = json.load(f)
                    tracks = data.get("tracks", [])
                    if not tracks and isinstance(data, list):
                        tracks = data
                    logger.info(f"Загружено {len(tracks)} треков")
                    return tracks
                except Exception as e:
                    logger.error(f"Ошибка загрузки треков из {path}: {e}")
        logger.error("Файлы с треками не найдены")
        return []

    def _process_single_audio_segment(self, task_data):
        rels_path, track, slide_num, media_dir = task_data
        try:
            logger.info(f"Обработка аудио для слайда {slide_num}")
            tree = ET.parse(rels_path)
            root = tree.getroot()
            targets = [
                rel.attrib.get("Target", "")
                for rel in root.findall("{http://schemas.openxmlformats.org/package/2006/relationships}Relationship")
                if rel.attrib.get("Target", "").lower().endswith(".mp3")
            ]
            if not targets:
                logger.warning(f"Не найдены MP3 цели для слайда {slide_num}")
                return None

            track_path = track.get("file_path") or track.get("path", "")
            real_path = self._find_track_file_fast(track_path)
            if not real_path:
                return None

            audio = AudioSegment.from_file(real_path)
            total_audio_length = len(audio)
            seg_start = int(float(track.get("segment_start", 0)) * 1000)
            raw_duration = float(track.get("segment_duration", self.default_ms / 1000))
            seg_dur = self._validate_segment_duration(raw_duration)
            end_ms = min(total_audio_length, seg_start + seg_dur)
            if seg_start > 0 or end_ms < total_audio_length:
                clip = audio[seg_start:end_ms]
            else:
                clip = audio

            clip = self._normalize_audio_loudness(clip)
            clip = self._apply_fade_effects(clip)

            out_media_path = media_dir / os.path.basename(targets[0])
            clip.export(out_media_path, format="mp3", bitrate="96k")
            logger.info(f"Аудио для слайда {slide_num} готово")
            return (slide_num, track)

        except Exception as e:
            logger.error(f"Ошибка обработки аудио для слайда {slide_num}: {e}")
            return None

    def _process_audio_parallel(self, audio_tasks):
        num_workers = mp.cpu_count()
        logger.info(f"Параллельная обработка {len(audio_tasks)} аудио задач на {num_workers} ядрах")
        results = []
        with ThreadPoolExecutor(max_workers=num_workers) as executor:
            future_to_task = {executor.submit(self._process_single_audio_segment, task): task for task in audio_tasks}
            completed = 0
            for future in as_completed(future_to_task):
                try:
                    result = future.result()
                    if result:
                        results.append(result)
                    completed += 1
                    logger.info(f"Прогресс аудио: {completed}/{len(audio_tasks)}")
                except Exception as e:
                    logger.error(f"Ошибка в потоке: {e}")
                    completed += 1
        logger.info(f"Все аудио файлы обработаны: {len(results)} успешно")
        return results

    def _add_track_number(self, slide, slide_width, slide_height, round_num):
        try:
            slide_num = int(''.join(ch for ch in slide.part.partname if ch.isdigit()) or 0)
            
            if round_num in self._round_track_counters:
                self._round_track_counters[round_num] += 1
                track_number = self._round_track_counters[round_num]
            else:
                track_number = ((slide_num - 5) % 40) + 1
            
            left = Emu(20272375)
            top = Emu(274320)
            width = Emu(914400)
            height = Emu(914400)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            
            run = p.add_run()
            run.text = str(track_number)
            font = run.font
            font.name = "Montserrat SemiBold"
            font.size = Pt(100)
            font.bold = True
            font.color.rgb = RGBColor(255, 0, 0)
            
            logger.debug(f"Добавлен номер трека {track_number} на слайд {slide_num} (раунд {round_num})")
            return True
            
        except Exception as e:
            logger.error(f"Ошибка добавления номера трека: {e}")
            return False

    def _apply_text_formatting_fast(self, original_run, new_text, is_artist, slide_num, artist_name=None):
        try:
            RED = RGBColor(255, 0, 0)
            BLACK = RGBColor(0, 0, 0)
            font_name = "Montserrat SemiBold"
            font_size = Pt(100)
            font_bold = True
            special_chars = {'@', '#', '$', '&', '€', '£', '¥'}
            special_patterns = {'zz', 'xxx', 'www', 'qq', 'kk'}
            words = new_text.split()
            if slide_num not in self._red_words_per_slide:
                self._red_words_per_slide[slide_num] = 0
            cached_red_words = set()
            if is_artist and artist_name and artist_name in self._artist_red_words_cache:
                cached_red_words = self._artist_red_words_cache[artist_name]
            paintable_words = []
            for word in words:
                if word.lower() in cached_red_words:
                    paintable_words.append((word, True))
                    continue
                has_special_char = any(char in word for char in special_chars)
                has_special_pattern = any(pattern in word.lower() for pattern in special_patterns)
                paintable_words.append((word, has_special_char or has_special_pattern))
            if (not any(paintable for _, paintable in paintable_words) and 
                words and self._red_words_per_slide[slide_num] == 0):
                random_index = random.randint(0, len(words) - 1)
                paintable_words[random_index] = (words[random_index], True)
            red_words_for_artist = set()
            for word, should_paint in paintable_words:
                if should_paint:
                    red_words_for_artist.add(word.lower())
            if is_artist and artist_name and red_words_for_artist:
                self._artist_red_words_cache[artist_name] = red_words_for_artist
            original_run.text = ""
            parent_paragraph = original_run._parent
            for word, should_paint in paintable_words:
                can_paint_red = should_paint and self._red_words_per_slide[slide_num] < 1
                if can_paint_red:
                    self._red_words_per_slide[slide_num] += 1
                    for char in word:
                        new_run = parent_paragraph.add_run()
                        new_run.text = char
                        new_run.font.name = font_name
                        new_run.font.size = font_size
                        new_run.font.bold = font_bold
                        if char in special_chars or (should_paint and not any(c in word for c in special_chars)):
                            new_run.font.color.rgb = RED
                        else:
                            new_run.font.color.rgb = BLACK
                    space_run = parent_paragraph.add_run()
                    space_run.text = " "
                    space_run.font.name = font_name
                    space_run.font.size = font_size
                    space_run.font.bold = font_bold
                    space_run.font.color.rgb = BLACK
                else:
                    new_run = parent_paragraph.add_run()
                    new_run.text = word + " "
                    new_run.font.name = font_name
                    new_run.font.size = font_size
                    new_run.font.bold = font_bold
                    new_run.font.color.rgb = BLACK
        except Exception as e:
            logger.error(f"Ошибка форматирования текста: {e}")
            original_run.text = new_text

    def _move_button_to_corner(self, slide, slide_height, slide_width, mirror=False):
        try:
            for shape in slide.shapes:
                if hasattr(shape, 'image'):
                    try:
                        if hasattr(shape, '_element'):
                            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                            blip_elements = shape._element.findall('.//a:blip', ns)
                            if blip_elements:
                                blip = blip_elements[0]
                                rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                                if rId and rId in slide.part.rels:
                                    rel = slide.part.rels[rId]
                                    image_part = rel.target_part
                                    if hasattr(image_part, 'partname'):
                                        filename = image_part.partname.split('/')[-1]
                                        if filename.lower() == 'image42.png':
                                            if mirror:
                                                new_left = Inches(0)
                                                new_top = slide_height - shape.height
                                                shape.left = new_left
                                                shape.top = new_top
                                                original_width = shape.width
                                                shape.width = -original_width
                                                shape.left = new_left + original_width
                                            else:
                                                shape.left = slide_width - shape.width
                                                shape.top = slide_height - shape.height
                                            return True
                    except Exception:
                        continue
            return False
        except Exception as e:
            logger.error(f"Ошибка перемещения кнопки: {e}")
            return False

    def _insert_artist_photo_fast(self, slide, processed_img, slide_height, slide_width, template_type):
        try:
            if processed_img is None:
                return
            temp_img_path = Path(tempfile.gettempdir()) / f"temp_photo_{random.randint(1000,9999)}.png"
            processed_img.save(temp_img_path, "PNG", optimize=True)
            img_width_px, img_height_px = processed_img.size
        
            dpi = 96
            width_inches = img_width_px / dpi
            height_inches = img_height_px / dpi
        
            if template_type == 1:
                left = Inches(0)
                top = slide_height - Inches(height_inches)
            elif template_type == 2:
                left = slide_width - Inches(width_inches)
                top = slide_height - Inches(height_inches)
            elif template_type == 3:
                left = Inches(0)
                top = slide_height - Inches(height_inches)
        
            slide.shapes.add_picture(str(temp_img_path), left, top, width=Inches(width_inches), height=Inches(height_inches))
        
            try:
                temp_img_path.unlink(missing_ok=True)
            except:
                pass
        
            logger.info(f"Фото добавлено: {'горизонтальное' if template_type == 3 else 'вертикальное'}")
        
        except Exception as e:
            logger.error(f"Ошибка вставки фото: {e}")

    def _replace_placeholders_and_photos_fast(self, prs, slide_track_map, make_bw: bool, rounds_config: list, rounds_count: int):
        slide_height = prs.slide_height
        slide_width = prs.slide_width
        logger.info(f"Размеры слайда: {slide_width} x {slide_height}")
        logger.info(f"Обработка {len(slide_track_map)} слайдов")
        
        self._red_words_per_slide.clear()
        self._round_track_counters = {1: 0, 2: 0, 3: 0}
        
        slide_to_round_map = {}
        for slide_num in slide_track_map.keys():
            if 5 <= slide_num <= 44:
                slide_to_round_map[slide_num] = 1
            elif 48 <= slide_num <= 87:
                slide_to_round_map[slide_num] = 2
            elif 91 <= slide_num <= 130:
                slide_to_round_map[slide_num] = 3
            else:
                slide_to_round_map[slide_num] = 1
        
        processed_slides = 0
        
        for i, slide in enumerate(prs.slides):
            slide_num = int(''.join(ch for ch in slide.part.partname if ch.isdigit()) or 0)
            if slide_num not in slide_track_map:
                continue

            track = slide_track_map[slide_num]
            artist = track.get("artist", "Неизвестный исполнитель")
            title = track.get("title", "Без названия")
            title = title.upper()
            
            logger.info(f"Обработка слайда {slide_num}: {artist} - {title}")
            
            current_round = slide_to_round_map.get(slide_num, 1)
            self._add_track_number(slide, slide_width, slide_height, current_round)
            
            image_path = track.get("image_path", "")
            template_type = 1
            if image_path:
                orientation = self._get_image_orientation(image_path)
                if orientation == "vertical":
                    template_type = random.choice([1, 2])
                else:
                    template_type = 3
                processed_img = self._load_and_process_image_fast(image_path, make_bw)
                if processed_img:
                    self._insert_artist_photo_fast(slide, processed_img, slide_height, slide_width, template_type)
                    logger.info(f"Фото добавлено на слайд {slide_num}")

            artist_shapes = []
            title_shapes = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "{{ARTIST}}" in run.text:
                            artist_shapes.append((shape, run))
                        elif "{{TRACK}}" in run.text:
                            title_shapes.append((shape, run))

            if artist_shapes:
                artist_shape, artist_run = artist_shapes[0]
                text_height = Inches(2.8)
                
                artist_needs_wrap = (len(artist.split()) > 1) and (len(artist) > 16)
                title_needs_wrap = (len(title.split()) > 1) and (len(title) > 16)
                use_word_wrap = artist_needs_wrap or title_needs_wrap

                if template_type == 1:
                    artist_shape.left = slide_width - Inches(12.0)
                    artist_shape.top = Inches(2.3)
                    artist_shape.width = Inches(12.0)
                    artist_shape.height = text_height
                    title = f'"{title}"' if title_shapes else title
                    text_alignment = PP_ALIGN.RIGHT
                elif template_type == 2:
                    artist_shape.left = Inches(0.5)
                    artist_shape.top = Inches(2.3)
                    artist_shape.width = Inches(12.0)
                    artist_shape.height = text_height
                    title = f'"{title}"' if title_shapes else title
                    text_alignment = PP_ALIGN.LEFT
                elif template_type == 3:
                    artist_shape.left = Inches(0.5)
                    artist_shape.top = Inches(0.2)
                    artist_shape.width = Inches(18.0)
                    artist_shape.height = text_height
                    title = f'"{title}"' if title_shapes else title
                    text_alignment = PP_ALIGN.LEFT

                artist_shape.text_frame.word_wrap = use_word_wrap
                artist_shape.text_frame.auto_size = None
                
                for paragraph in artist_shape.text_frame.paragraphs:
                    paragraph.alignment = text_alignment
                
                self._apply_text_formatting_fast(artist_run, artist, True, slide_num, artist)
                logger.info(f"Текст артиста добавлен: {artist}")

                if title_shapes:
                    title_shape, title_run = title_shapes[0]
                    if template_type == 1:
                        title_shape.left = slide_width - Inches(12.0)
                        title_shape.top = Inches(4.3)
                        title_shape.width = Inches(12.0)
                        title_shape.height = text_height
                    elif template_type == 2:
                        title_shape.left = Inches(0.5)
                        title_shape.top = Inches(4.3)
                        title_shape.width = Inches(12.0)
                        title_shape.height = text_height
                    elif template_type == 3:
                        title_shape.left = Inches(0.5)
                        title_shape.top = Inches(2.0)
                        title_shape.width = Inches(18.0)
                        title_shape.height = text_height

                    title_shape.text_frame.word_wrap = use_word_wrap
                    title_shape.text_frame.auto_size = None
                    
                    for paragraph in title_shape.text_frame.paragraphs:
                        paragraph.alignment = text_alignment
                    
                    self._apply_text_formatting_fast(title_run, title, False, slide_num, artist)
                    logger.info(f"Текст трека добавлен: {title}")

            if template_type == 2:
                self._move_button_to_corner(slide, slide_height, slide_width, mirror=True)
            else:
                self._move_button_to_corner(slide, slide_height, slide_width, mirror=False)

            processed_slides += 1
            logger.info(f"Слайд {slide_num} обработан ({processed_slides}/{len(slide_track_map)}) - Раунд {current_round}")

        logger.info(f"Все слайды обработаны: {processed_slides} из {len(slide_track_map)}")
        
        for round_num in range(1, rounds_count + 1):
            if round_num in self._round_track_counters:
                logger.info(f"Раунд {round_num}: {self._round_track_counters[round_num]} треков")

    def generate(self, game_title: str, tracks: list = None, make_bw: bool = False, 
                 use_parallel: bool = True, segment_duration: int = None,
                 rounds_config: list = None, rounds_count: int = 1):
        
        logger.info(f"Запуск генерации презентации с {rounds_count} раундами")
        
        if segment_duration is not None:
            validated_ms = self._validate_segment_duration(segment_duration)
            self.default_ms = validated_ms
            logger.info(f"Кастомная длительность: {segment_duration} сек")
            
        if not tracks:
            tracks = self._load_tracks_from_json_fast()
        if not tracks:
            raise ValueError("Нет треков для генерации")
            
        if not rounds_config:
            default_tracks = min(40, len(tracks))
            rounds_config = [default_tracks]
            logger.info(f"Используем дефолтную конфиг: {rounds_config}")
        else:
            rounds_config = rounds_config[:rounds_count]
            logger.info(f"Конфигурация раундов: {rounds_config}")
        
        total_tracks_needed = sum(rounds_config)
        if len(tracks) < total_tracks_needed:
            raise ValueError(f"Недостаточно треков: нужно {total_tracks_needed}, есть {len(tracks)}")
        
        logger.info(f"Треков для обработки: {total_tracks_needed}")
        
        slides_to_delete = self._calculate_slides_to_delete(rounds_config, rounds_count)
        
        tmp = Path(tempfile.mkdtemp(prefix="pptx_fast_"))
        extract_dir = tmp / "extracted"
        extract_dir.mkdir(parents=True, exist_ok=True)
        
        try:
            logger.info("Распаковка шаблона...")
            with zipfile.ZipFile(self.base_path, "r") as z:
                z.extractall(extract_dir)
                
            slides_dir = extract_dir / "ppt" / "slides"
            slides_rels_dir = slides_dir / "_rels"
            media_dir = extract_dir / "ppt" / "media"
            media_dir.mkdir(parents=True, exist_ok=True)
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            out_root = Path.cwd() / "output" / f"presentation_{timestamp}"
            out_root.mkdir(parents=True, exist_ok=True)
            
            slide1 = slides_dir / "slide1.xml"
            if slide1.exists():
                content = slide1.read_text(encoding="utf-8")
                if "{{TITLE}}" in content:
                    slide1.write_text(content.replace("{{TITLE}}", game_title), encoding="utf-8")
                    logger.info("Заголовок презентации обновлен")
                    
            rels_files = sorted(
                [f for f in slides_rels_dir.glob("slide*.xml.rels")],
                key=lambda x: int(''.join(filter(str.isdigit, x.stem)) or 0)
            )
            
            slide_track_map = {}
            audio_tasks = []
            current_track_index = 0
            
            for round_idx in range(rounds_count):
                round_num = round_idx + 1
                tracks_in_round = rounds_config[round_idx]
                start_slide, _ = self.round_slide_ranges[round_num]
                
                for i in range(tracks_in_round):
                    if current_track_index >= len(tracks):
                        break
                        
                    track = tracks[current_track_index]
                    slide_num = start_slide + i
                    
                    rels_path = slides_rels_dir / f"slide{slide_num}.xml.rels"
                    if rels_path.exists():
                        audio_tasks.append((rels_path, track, slide_num, media_dir))
                        slide_track_map[slide_num] = track
                    
                    current_track_index += 1
            
            logger.info(f"Распределено {current_track_index} треков по {len(slide_track_map)} слайдам")
            
            logger.info("Параллельная обработка аудио...")
            if use_parallel and len(audio_tasks) > 1:
                results = self._process_audio_parallel(audio_tasks)
            else:
                results = []
                for task in audio_tasks:
                    result = self._process_single_audio_segment(task)
                    if result:
                        results.append(result)
                        
            for slide_num, processed_track in results:
                slide_track_map[slide_num] = processed_track
                
            final_pptx = out_root / f"presentation_{timestamp}.pptx"
            
            logger.info("Создание PPTX...")
            with zipfile.ZipFile(final_pptx, "w", zipfile.ZIP_DEFLATED, compresslevel=6) as zip_out:
                for root, _, files in os.walk(extract_dir):
                    for file in files:
                        fp = Path(root) / file
                        arcname = os.path.relpath(fp, extract_dir)
                        zip_out.write(fp, arcname)
            
            logger.info("Замена текста и фото...")
            prs = Presentation(final_pptx)
            self._replace_placeholders_and_photos_fast(prs, slide_track_map, make_bw, rounds_config, rounds_count)
            prs.save(final_pptx)
            
            if slides_to_delete:
                logger.info(f"Удаление {len(slides_to_delete)} слайдов...")
                self._delete_slides_directly(final_pptx, slides_to_delete)
                logger.info(f"Слайды удалены! Режим: {rounds_count} раунда(ов)")
            else:
                logger.info("Нет слайдов для удаления")
            
            logger.info(f"Презентация готова: {final_pptx}")
            return str(out_root)
            
        except Exception as e:
            logger.error(f"Ошибка генерации: {e}")
            import traceback
            logger.error(traceback.format_exc())
            raise
        finally:
            try:
                shutil.rmtree(tmp, ignore_errors=True)
                logger.info("Временные файлы очищены")
            except:
                pass


if __name__ == "__main__":
    try:
        logger.info("Запуск генератора с раундами")
        generator = ModernPresentationGenerator("template.pptx")
        
        result_path = generator.generate(
            game_title="Моя викторина 1 раунд",
            tracks=None,
            make_bw=False,
            use_parallel=True,
            rounds_config=[20],
            rounds_count=1
        )
        
        print(f"Презентация создана: {result_path}")
        
    except Exception as e:
        print(f"Ошибка: {e}")
        import traceback
        print(traceback.format_exc())